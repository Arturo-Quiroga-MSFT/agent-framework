#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Microsoft Agent Framework Overview.

This script creates a comprehensive presentation covering the Agent Framework's
architecture, capabilities, and getting started guide.

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def create_title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide


def create_content_slide(prs, title, content_items, layout_index=1):
    """Create a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add content
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    return slide


def create_two_column_slide(prs, title, left_items, right_items):
    """Create a slide with two columns of content."""
    slide_layout = prs.slide_layouts[3]  # Two content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Left column
    left_shape = slide.placeholders[1]
    left_frame = left_shape.text_frame
    left_frame.clear()
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    # Right column
    right_shape = slide.placeholders[2]
    right_frame = right_shape.text_frame
    right_frame.clear()
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    return slide


def create_code_slide(prs, title, code_text, language="python"):
    """Create a slide with code example."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add code box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    code_box = slide.shapes.add_textbox(left, top, width, height)
    code_frame = code_box.text_frame
    code_frame.text = code_text
    code_frame.paragraphs[0].font.name = 'Courier New'
    code_frame.paragraphs[0].font.size = Pt(14)
    code_frame.word_wrap = True
    
    # Add background color to code box
    fill = code_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide


def main():
    """Generate the Agent Framework overview presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Microsoft Agent Framework",
        "Building, Orchestrating, and Deploying AI Agents\nPrivate Preview Overview"
    )
    
    # Slide 2: What is Agent Framework?
    create_content_slide(
        prs,
        "What is Microsoft Agent Framework?",
        [
            "🤖 Comprehensive multi-language framework for AI agents",
            "🔄 Support for both .NET and Python implementations",
            "🎯 Everything from simple chat agents to complex workflows",
            "🌐 Graph-based orchestration for multi-agent systems",
            "🔌 Extensible plugin ecosystem",
            "☁️ Enterprise-ready with Azure integration"
        ]
    )
    
    # Slide 3: Current Status
    create_content_slide(
        prs,
        "Current Status - Private Preview",
        [
            "📦 Status: Active Private Preview",
            "🔄 Packages: Nightly builds available via GitHub Packages",
            "🚀 Ready for: Testing, feedback, and early adoption",
            "📝 Not yet: Public PyPI or NuGet releases",
            "💬 Feedback: GitHub issues and surveys welcomed",
            "⚡ Updates: Active development, sync regularly"
        ]
    )
    
    # Slide 4: Key Features
    create_content_slide(
        prs,
        "Key Features & Capabilities",
        [
            "🎭 Multi-Agent Orchestration: Group chat, sequential, concurrent patterns",
            "📊 Graph-based Workflows: Data flows with streaming & checkpointing",
            "🔧 Plugin Ecosystem: Native functions, OpenAPI, Model Context Protocol",
            "🧠 LLM Support: OpenAI, Azure OpenAI, Azure AI Foundry",
            "⚙️ Runtime Support: In-process and distributed agent execution",
            "🎨 Multimodal: Text, vision, and function calling",
            "🔄 Human-in-the-Loop: Interactive approvals and interventions"
        ]
    )
    
    # Slide 5: Architecture Overview
    create_content_slide(
        prs,
        "Architecture Components",
        [
            "1. Chat Agents - Conversational AI agents with instructions, tools, and state management for focused tasks",
            "2. Agent Executors - Wrapper layer that integrates agents into workflows, enabling coordination between multiple agents",
            "3. Workflow Builder - Fluent API for creating graph-based orchestrations with edges, streaming, and checkpointing",
            "4. Chat Clients - Provider-agnostic abstractions supporting OpenAI, Azure OpenAI, and Azure AI Foundry",
            "5. Tools & Plugins - Type-safe function definitions with automatic schema generation for LLM invocation",
            "6. Runtime - Flexible execution supporting both in-process agents and distributed multi-agent systems"
        ]
    )
    
    # Slide 6: Multi-Language Support
    create_two_column_slide(
        prs,
        "Cross-Platform Support",
        [
            "🐍 Python",
            "• Version: 3.10+",
            "• Platforms: Windows, macOS, Linux",
            "• Package: agent-framework",
            "• Extensions: azure-ai, microsoft",
            "• Dev Tools: uv, pip, venv"
        ],
        [
            ".NET",
            "• Versions: .NET 9.0, 8.0, Framework 4.7.2",
            "• Platforms: Windows, macOS, Linux",
            "• Package: Microsoft.Agents.AI",
            "• Extensions: Azure, Abstractions",
            "• Dev Tools: dotnet CLI, NuGet"
        ]
    )
    
    # Slide 7: Agent Types
    create_content_slide(
        prs,
        "Types of Agents",
        [
            "💬 ChatAgent - Simple conversational agents with tools",
            "🔄 AgentExecutor - Agents wrapped for workflows",
            "🎯 Custom Executors - Python/C# code executors in workflows",
            "🌐 Azure AI Agents - Managed agents with Azure AI Foundry",
            "🤖 Azure OpenAI Assistants - Stateful assistants",
            "📡 Copilot Studio Agents - Microsoft 365 integration"
        ]
    )
    
    # Slide 8: Orchestration Patterns
    create_content_slide(
        prs,
        "Multi-Agent Orchestration Patterns",
        [
            "🔗 Sequential - Agents execute in order",
            "⚡ Concurrent (Fan-out/Fan-in) - Parallel execution",
            "🎯 Handoff - Transfer control between agents",
            "💬 Group Chat - Multi-agent conversations",
            "🌳 Hierarchical - Leader-worker patterns",
            "🔄 Conditional Routing - Dynamic agent selection"
        ]
    )
    
    # Slide 9: Workflow Capabilities
    create_content_slide(
        prs,
        "Graph-Based Workflows",
        [
            "📊 Visual Workflow Builder - Fluent API for graph construction",
            "🎬 Streaming Support - Real-time event streaming",
            "💾 Checkpointing - Save and restore workflow state",
            "⏮️ Time-Travel Debugging - Replay workflow execution",
            "🤝 Human-in-the-Loop - Interactive approvals",
            "📈 Visualization - Mermaid and GraphViz diagrams"
        ]
    )
    
    # Slide 10: Python Code Example
    python_code = '''# Create a weather agent with tools
from agent_framework import ChatAgent
from agent_framework.azure import AzureOpenAIChatClient

def get_weather(location: str) -> str:
    """Get weather for a location."""
    return f"Weather in {location}: Sunny, 72°F"

agent = ChatAgent(
    name="WeatherAssistant",
    instructions="Help users with weather info",
    chat_client=AzureOpenAIChatClient(),
    tools=[get_weather]
)

response = await agent.run("What's the weather in Seattle?")
print(response.text)'''
    
    create_code_slide(prs, "Python Example - Simple Agent", python_code)
    
    # Slide 11: .NET Code Example
    dotnet_code = '''// Create a weather agent with tools
using Microsoft.Agents.AI;
using Microsoft.Agents.AI.Azure;

string GetWeather(string location) =>
    $"Weather in {location}: Sunny, 72°F";

var agent = new AzureOpenAIClient(
        new Uri(Environment.GetEnvironmentVariable("AZURE_OPENAI_ENDPOINT")!),
        new AzureCliCredential())
    .GetChatClient("gpt-4o")
    .CreateAIAgent(
        instructions: "Help users with weather information",
        tools: [AIFunctionFactory.Create(GetWeather)]);

var response = await agent.RunAsync("What's the weather in Seattle?");
Console.WriteLine(response.Text);'''
    
    create_code_slide(prs, ".NET Example - Simple Agent", dotnet_code)
    
    # Slide 12: Workflow Example
    workflow_code = '''# Create a multi-agent workflow
from agent_framework import WorkflowBuilder, AgentExecutor

# Create agents
researcher = AgentExecutor(chat_client.create_agent(
    instructions="Research facts"))
writer = AgentExecutor(chat_client.create_agent(
    instructions="Write content"))

# Build workflow
workflow = (WorkflowBuilder()
    .set_start_executor(researcher)
    .add_edge(researcher, writer)
    .build())

# Run workflow
async for event in workflow.run_stream("Write about AI"):
    print(event)'''
    
    create_code_slide(prs, "Workflow Example - Sequential Agents", workflow_code)
    
    # Slide 13: Plugin Ecosystem
    create_content_slide(
        prs,
        "Plugin & Tool Ecosystem",
        [
            "⚡ Native Functions - Python/C# functions with annotations",
            "🌐 OpenAPI - REST API integrations (planned)",
            "🔌 Model Context Protocol (MCP) - Tool interoperability",
            "🎯 Azure AI Search - Knowledge base integration (planned)",
            "🔍 Bing Grounding - Web search capabilities (planned)",
            "🔧 Custom Tools - Build your own extensions"
        ]
    )
    
    # Slide 14: Azure Integration
    create_content_slide(
        prs,
        "Azure Integration",
        [
            "🎯 Azure AI Foundry - Managed agent infrastructure",
            "🧠 Azure OpenAI - GPT models and assistants",
            "🔐 Azure Identity - AAD authentication",
            "💾 Azure Cosmos DB - Workflow state storage",
            "📊 Azure Monitor - Telemetry and logging",
            "🔍 Application Insights - Performance tracking"
        ]
    )
    
    # Slide 15: Getting Started - Installation
    create_content_slide(
        prs,
        "Getting Started - Installation",
        [
            "Python:",
            "  pip install agent-framework",
            "  pip install agent-framework[azure-ai,microsoft,all]",
            "",
            ".NET:",
            "  dotnet add package Microsoft.Agents.AI --version 0.0.1-nightly-*",
            "",
            "Or run samples directly from the cloned repository!"
        ]
    )
    
    # Slide 16: Configuration
    create_content_slide(
        prs,
        "Configuration - Environment Variables",
        [
            "OpenAI:",
            "  OPENAI_API_KEY, OPENAI_CHAT_MODEL_ID",
            "",
            "Azure OpenAI:",
            "  AZURE_OPENAI_API_KEY, AZURE_OPENAI_ENDPOINT",
            "  AZURE_OPENAI_CHAT_DEPLOYMENT_NAME",
            "",
            "Azure AI Foundry:",
            "  AZURE_AI_PROJECT_ENDPOINT, AZURE_AI_MODEL_DEPLOYMENT_NAME"
        ]
    )
    
    # Slide 17: Sample Projects
    create_two_column_slide(
        prs,
        "Sample Projects & Examples",
        [
            "Python Samples:",
            "• Basic agents with tools",
            "• Azure integration",
            "• Workflow orchestration",
            "• Multi-agent patterns",
            "• Real weather API demo",
            "• Notebook tutorials"
        ],
        [
            ".NET Samples:",
            "• Minimal console agent",
            "• Agent providers",
            "• Orchestration patterns",
            "• Semantic Kernel migration",
            "• A2A client/server",
            "• Web chat UI"
        ]
    )
    
    # Slide 18: DevUI and Debugging
    create_content_slide(
        prs,
        "Development Tools",
        [
            "🖥️ DevUI - Web-based agent testing interface",
            "🐛 Debugging - Step through agent execution",
            "📊 Visualization - WorkflowViz for graph diagrams",
            "📝 Logging - OpenTelemetry instrumentation",
            "🔍 Tracing - Track agent interactions",
            "⏮️ Time-Travel - Replay and debug workflows"
        ]
    )
    
    # Slide 19: Best Practices
    create_content_slide(
        prs,
        "Best Practices",
        [
            "🔐 Security - Use Azure Identity, avoid hardcoded keys",
            "📝 Instructions - Clear, specific agent instructions",
            "🛠️ Tools - Type hints and descriptions for all parameters",
            "🔄 Error Handling - Graceful failures and retries",
            "📊 Monitoring - Implement telemetry and logging",
            "🧪 Testing - Unit test agents and workflows",
            "📚 Documentation - Document agent capabilities"
        ]
    )
    
    # Slide 20: Resources
    create_content_slide(
        prs,
        "Resources & Documentation",
        [
            "📚 Repository: github.com/microsoft/agent-framework",
            "📖 Python Docs: user-documentation-python/",
            "📖 .NET Docs: user-documentation-dotnet/",
            "🏗️ Architecture Decisions: docs/decisions/",
            "🎨 Design Docs: docs/design/",
            "🐛 Issues: GitHub Issues",
            "💬 Feedback: Microsoft Forms survey"
        ]
    )
    
    # Slide 21: Roadmap & Future
    create_content_slide(
        prs,
        "Current Focus & Evolution",
        [
            "🚀 Active Development Areas:",
            "• Completing core agent patterns",
            "• Enhanced workflow capabilities",
            "• Additional Azure service integrations",
            "• OpenAPI and MCP tool support",
            "• Performance optimizations",
            "• Documentation and samples",
            "",
            "📦 Moving toward public preview and GA releases"
        ]
    )
    
    # Slide 22: Call to Action
    create_content_slide(
        prs,
        "Get Involved!",
        [
            "🔧 Try the Framework:",
            "  Clone the repo and run samples",
            "",
            "💬 Provide Feedback:",
            "  File issues, complete surveys",
            "",
            "🤝 Contribute:",
            "  Review contribution guidelines",
            "",
            "📢 Stay Updated:",
            "  Sync your repo regularly for latest features"
        ]
    )
    
    # Slide 23: Q&A
    create_title_slide(
        prs,
        "Questions?",
        "Microsoft Agent Framework\n\nThank you!"
    )
    
    # Save presentation
    output_file = "agent_framework_overview.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"\n🎯 Next steps:")
    print(f"   1. Open the file in PowerPoint or compatible software")
    print(f"   2. Customize colors, fonts, and layouts as needed")
    print(f"   3. Add company branding if required")
    print(f"   4. Review and practice before presenting")


if __name__ == "__main__":
    main()
